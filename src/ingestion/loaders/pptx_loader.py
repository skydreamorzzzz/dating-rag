"""Loader for PPTX files."""

from pathlib import Path

from src.ingestion.loaders.base import BaseLoader
from src.ingestion.schemas import Document
from src.utils.logger import get_logger

logger = get_logger(__name__)


class PptxLoader(BaseLoader):
    """Load a .pptx file, returning one Document per slide."""

    def load(self, file_path: Path) -> list[Document]:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise RuntimeError(
                "python-pptx is required to load PPTX files. Install requirements.txt first."
            ) from e

        try:
            presentation = Presentation(str(file_path))
        except Exception as e:
            raise ValueError(f"Failed to open PPTX {file_path.name}: {e}") from e

        docs: list[Document] = []
        for slide_no, slide in enumerate(presentation.slides, start=1):
            parts = self._extract_slide_text(slide)
            text = "\n".join(part for part in parts if part).strip()
            warnings: list[str] = []
            if not text:
                warnings.append("No text extracted from PPTX slide — slide may be image-only")

            docs.append(
                Document(
                    source_file=file_path.name,
                    source_path=str(file_path),
                    doc_type="pptx",
                    text=text,
                    page=slide_no,
                    section=f"slide {slide_no}",
                    warnings=warnings,
                )
            )

        if not docs:
            docs.append(
                Document(
                    source_file=file_path.name,
                    source_path=str(file_path),
                    doc_type="pptx",
                    text="",
                    warnings=["No slides found in PPTX"],
                )
            )

        total_chars = sum(len(doc.text) for doc in docs)
        logger.info("Loaded %s (%d slides, %d chars)", file_path.name, len(docs), total_chars)
        return docs

    def _extract_slide_text(self, slide) -> list[str]:
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    parts.append(text)

            if getattr(shape, "has_table", False):
                rows: list[str] = []
                for row in shape.table.rows:
                    rows.append("\t".join(cell.text for cell in row.cells))
                if rows:
                    parts.append("\n".join(rows))
        return parts

    @staticmethod
    def supported_extensions() -> list[str]:
        return [".pptx"]
